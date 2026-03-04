"""PptxExtractor — extracts text from PPTX files using python-pptx."""

from __future__ import annotations

from typing import Any, ClassVar

from ..core.base import BaseExtractor, ExtractionError, ExtractionResult


class PptxExtractor(BaseExtractor):
    """Extracts text from PowerPoint PPTX files using python-pptx."""

    supported_extensions: ClassVar[set[str]] = {".pptx"}
    required_packages: ClassVar[set[str]] = {"pptx"}

    def extract(self, source: str) -> ExtractionResult:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ExtractionError(
                "python-pptx is required for PPTX extraction: pip install python-pptx",
                source=source,
                cause=e,
            )

        try:
            prs = Presentation(source)
        except Exception as e:
            raise ExtractionError(
                f"Failed to open PPTX {source}: {e}", source=source, cause=e
            )

        slides_text: list[str] = []
        slides_md: list[str] = []
        metadata: dict[str, Any] = {}

        metadata["Slides"] = len(prs.slides)
        core = prs.core_properties
        if core.author:
            metadata["Author"] = core.author
        if core.title:
            metadata["Title"] = core.title

        for idx, slide in enumerate(prs.slides, 1):
            text_parts: list[str] = []
            md_parts: list[str] = [f"## Slide {idx}"]

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if not para_text:
                            continue
                        text_parts.append(para_text)
                        # Bullet indentation based on paragraph level
                        indent = "  " * para.level
                        md_parts.append(f"{indent}- {para_text}" if para.level > 0 else para_text)

                if shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    if table_md:
                        # Plain text: just the cell values
                        for row in shape.table.rows:
                            row_text = "\t".join(
                                cell.text.strip() for cell in row.cells
                            )
                            text_parts.append(row_text)
                        md_parts.append(table_md)

            if text_parts:
                slides_text.append("\n".join(text_parts))
            if len(md_parts) > 1:  # more than just the header
                slides_md.append("\n\n".join(md_parts))

        full_text = "\n\n".join(slides_text)
        markdown_text = "\n\n".join(slides_md)

        return ExtractionResult(
            text=full_text,
            source=source,
            source_type="pptx",
            extractor_name=self.__class__.__name__,
            metadata=metadata,
            markdown_text=markdown_text or None,
        )

    @staticmethod
    def _table_to_markdown(table: Any) -> str:
        """Render a PPTX table as markdown."""
        rows_data: list[list[str]] = []
        for row in table.rows:
            rows_data.append([cell.text.strip() for cell in row.cells])

        if not rows_data or len(rows_data) < 1:
            return ""

        cols = len(rows_data[0])
        header = rows_data[0]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join("---" for _ in range(cols)) + " |",
        ]
        for row in rows_data[1:]:
            padded = row + [""] * (cols - len(row))
            lines.append("| " + " | ".join(padded[:cols]) + " |")
        return "\n".join(lines)
